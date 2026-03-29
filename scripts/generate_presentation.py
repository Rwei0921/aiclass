from pptx import Presentation
from pptx.util import Inches, Pt

def add_slide(prs, layout, title, content):
    slide = prs.slides.add_slide(layout)
    if title:
        title_shape = slide.shapes.title
        title_shape.text = title
    if content and len(slide.placeholders) > 1:
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        for item in content:
            if isinstance(item, list):
                for subitem in item:
                    p = tf.add_paragraph()
                    p.text = subitem
                    p.level = 1
            else:
                p = tf.add_paragraph()
                p.text = item
                p.level = 0
    return slide

def create_midterm_presentation():
    prs = Presentation()
    
    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AI Explorer：互動式 AI 入門教學網站"
    subtitle.text = "期中成果報告\n\n「點 → 看 → 懂」的 AI 新手村"

    # Slide 2: 專案動機與目標
    bullet_slide_layout = prs.slide_layouts[1]
    add_slide(prs, bullet_slide_layout, "1. 專案動機與目標", [
        "核心動機：",
        ["解決一般人對 AI 的距離感與學習門檻", "將艱澀的 AI 知識轉化為視覺化、可互動的內容"],
        "專案目標：",
        ["打造一個「完全不懂 AI」也能輕鬆上手的教學網站", "使用者透過「點擊互動 + 短影片 + 圖解」來學習", "將學習過程遊戲化、探索化（AI 新手村概念）"]
    ])

    # Slide 3: 受眾分析
    add_slide(prs, bullet_slide_layout, "2. 受眾分析", [
        "目標對象（Target Audience）：",
        ["同學、家人、非資訊背景的社會大眾", "想了解 AI 卻不知從何開始的初學者"],
        "痛點分析：",
        ["名詞太多（機器學習、深度學習、神經網路太抽象）", "缺乏直觀的理解方式", "長篇大論的文字容易讓人放棄"],
        "我們的解決方案：",
        ["互動式可視化圖表 (AI Pipeline / 知識圖譜)", "每次只呈現一個核心概念 (白話解釋 + 短影片 + 測驗)"]
    ])

    # Slide 4: 網站整體架構 (Site Map)
    add_slide(prs, bullet_slide_layout, "3. 網站架構與互動設計", [
        "已完成的核心架構：",
        ["1. Home 首頁：引導使用者進入，展示核心理念", "2. Learn (Map) 互動學習：主頁，提供 AI 互動地圖", "3. Topics / Cases：技術主題清單與應用簡介", "4. About：團隊介紹與專案說明"],
        "核心互動設計 (Learn 頁面)：",
        ["左側：動態可視化地圖 (SVG)", "右側：解釋面板 (影片 + 圖解 + 白話說明)", "操作流程：點擊節點 → 面板立即更新 → 觀看與學習"]
    ])

    # Slide 5: 目前完成項目
    add_slide(prs, bullet_slide_layout, "4. 階段性成果 (目前完成進度)", [
        "基礎建設與資料層：",
        ["完成專案資料夾結構與開發設定", "建置 22 個 AI 概念節點的 JSON 資料庫"],
        "前端開發與設計 (路徑 A - 純前端實作)：",
        ["完成 CSS 全域樣式與地圖專用設計系統", "完成 index.html, map.html, topics.html 等核心頁面", "實作動態 SVG AI 地圖互動邏輯 (map.js)"],
        "內容層整合：",
        ["成功整合 NotebookLM 匯出的學習內容", "右側面板可動態播放 YouTube 影片並顯示說明"]
    ])

    # Slide 6: Demo 展示
    add_slide(prs, bullet_slide_layout, "5. 成果展示 (Demo)", [
        "展示重點：",
        ["1. 首頁進入地圖的引導流程", "2. 點擊 AI 節點 (例如：機器學習、深度學習)", "3. 觀察右側面板的「影片 + 說明」即時更新", "4. 展現流暢的 RWD 版面與 UI 回饋體驗"],
        "(切換至瀏覽器進行實際操作展示)"
    ])

    # Slide 7: 期末工作規劃
    add_slide(prs, bullet_slide_layout, "6. 期末工作規劃 (剩餘 50%)", [
        "內容與互動深化：",
        ["補齊所有 22 個節點的完整內容 (生活例子、常見誤解)", "每個節點加入「小測驗」與立即回饋機制", "加入「學習總測驗」單元，提供等級評估"],
        "進階功能與優化：",
        ["擴充 Cases 應用案例 (醫療/交通/教育等具體情境)", "深化 UI/UX (手機版 RWD 調整、動畫回饋)", "建立完整的 Report 反思與系統架構頁面"],
        "預計期末將能提供完整的產品級體驗！"
    ])

    prs.save(r'd:\AIClass\ppt\AI_Explorer_Midterm_Report.pptx')
    print("Midterm PPT generated successfully!")

if __name__ == '__main__':
    create_midterm_presentation()
