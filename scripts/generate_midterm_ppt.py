import json
import re
from collections import Counter
from datetime import datetime
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


COLORS = {
    "bg": RGBColor(7, 18, 38),
    "bg_soft": RGBColor(13, 34, 61),
    "primary": RGBColor(50, 197, 255),
    "accent": RGBColor(121, 240, 198),
    "text": RGBColor(236, 244, 255),
    "muted": RGBColor(157, 179, 207),
}


def style_slide_background(slide, title_text):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["bg"]

    top_band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.33),
        Inches(0.55),
    )
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = COLORS["bg_soft"]
    top_band.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0.47),
        Inches(13.33),
        Inches(0.08),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS["primary"]
    accent.line.fill.background()

    footer = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(7.35),
        Inches(13.33),
        Inches(0.15),
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = COLORS["bg_soft"]
    footer.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.08), Inches(8.8), Inches(0.35))
    title_tf = title_box.text_frame
    title_tf.clear()
    p = title_tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Noto Sans TC"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]

    tag_box = slide.shapes.add_textbox(Inches(10.2), Inches(0.08), Inches(2.6), Inches(0.3))
    tag_tf = tag_box.text_frame
    tag_tf.clear()
    p = tag_tf.paragraphs[0]
    p.text = "AI Explorer"
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = "Space Grotesk"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS["muted"]


def set_shape_text_style(shape, size=22, bold=False, color=None, font_name="Noto Sans TC"):
    tf = shape.text_frame
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color or COLORS["text"]


def style_content_title(shape):
    tf = shape.text_frame
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Space Grotesk"
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = COLORS["text"]


def add_page_number(slide, page_text):
    box = slide.shapes.add_textbox(Inches(12.2), Inches(7.08), Inches(0.8), Inches(0.24))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = page_text
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = "Space Grotesk"
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS["muted"]


def add_card(slide, x, y, w, h, title, points):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["bg_soft"]
    card.line.color.rgb = COLORS["primary"]
    card.line.width = Pt(1)

    title_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.12), w - Inches(0.4), Inches(0.5))
    tt = title_box.text_frame
    tt.clear()
    p = tt.paragraphs[0]
    p.text = title
    p.font.name = "Noto Sans TC"
    p.font.size = Pt(21)
    p.font.bold = True
    p.font.color.rgb = COLORS["accent"]

    body = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.72), w - Inches(0.35), h - Inches(0.85))
    bt = body.text_frame
    bt.clear()
    for i, point in enumerate(points):
        para = bt.paragraphs[0] if i == 0 else bt.add_paragraph()
        para.text = f"• {point}"
        para.level = 0
        para.font.name = "Noto Sans TC"
        para.font.size = Pt(17)
        para.font.color.rgb = COLORS["text"]


def ensure_visual_assets():
    asset_dir = Path("D:/AIClass/assets/images/ppt")
    asset_dir.mkdir(parents=True, exist_ok=True)

    specs = [
        ("map-overview.png", "互動地圖全覽", "時間軸 + 節點 + 右側面板", (16, 43, 77)),
        ("demo-turing.png", "圖靈測試節點", "NotebookLM 摘要 + 影片 + 來源", (20, 58, 93)),
        ("demo-transformer.png", "Transformer 節點", "模型架構與技術關聯展示", (24, 64, 102)),
        ("demo-chatgpt.png", "ChatGPT 節點", "應用突破與多模態趨勢", (18, 53, 88)),
    ]

    for filename, title, subtitle, bg in specs:
        path = asset_dir / filename
        if path.exists():
            continue

        img = Image.new("RGB", (1280, 720), bg)
        draw = ImageDraw.Draw(img)
        draw.rectangle((40, 40, 1240, 680), outline=(50, 197, 255), width=4)
        draw.rectangle((80, 90, 1200, 170), fill=(10, 25, 45))
        draw.rectangle((80, 210, 760, 640), fill=(13, 34, 61), outline=(121, 240, 198), width=3)
        draw.rectangle((790, 210, 1200, 640), fill=(9, 28, 50), outline=(50, 197, 255), width=3)
        draw.text((110, 110), title, fill=(236, 244, 255))
        draw.text((110, 150), subtitle, fill=(157, 179, 207))
        draw.text((115, 230), "AI Timeline Map", fill=(236, 244, 255))
        draw.text((810, 230), "Node Detail Panel", fill=(236, 244, 255))
        img.save(path)

    return {
        "overview": str(asset_dir / "map-overview.png"),
        "turing": str(asset_dir / "demo-turing.png"),
        "transformer": str(asset_dir / "demo-transformer.png"),
        "chatgpt": str(asset_dir / "demo-chatgpt.png"),
    }


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    style_slide_background(slide, title)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    style_content_title(slide.shapes.title)

    stf = slide.placeholders[1].text_frame
    for paragraph in stf.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Noto Sans TC"
            run.font.size = Pt(22)
            run.font.color.rgb = COLORS["muted"]

    hero = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(4.55),
        Inches(12.0),
        Inches(1.7),
    )
    hero.fill.solid()
    hero.fill.fore_color.rgb = COLORS["bg_soft"]
    hero.line.color.rgb = COLORS["primary"]
    hero.line.width = Pt(1.2)

    htf = hero.text_frame
    htf.clear()
    p = htf.paragraphs[0]
    p.text = "讓 AI 歷史變成可點擊、可理解、可驗證的學習地圖"
    p.font.name = "Noto Sans TC"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLORS["text"]
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    style_slide_background(slide, title)
    slide.shapes.title.text = title
    style_content_title(slide.shapes.title)
    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.75),
        Inches(1.7),
        Inches(11.9),
        Inches(4.95),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLORS["bg_soft"]
    panel.line.color.rgb = COLORS["primary"]
    panel.line.width = Pt(1)

    body_box = slide.shapes.add_textbox(Inches(1.05), Inches(2.0), Inches(11.2), Inches(4.4))
    body = body_box.text_frame
    body.clear()

    for idx, item in enumerate(bullets):
        p = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        p.font.name = "Noto Sans TC"
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS["text"]

    return slide


def add_image_bullets_slide(prs, title, image_path, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    style_slide_background(slide, title)
    slide.shapes.title.text = title
    style_content_title(slide.shapes.title)

    image_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(1.75), Inches(7.0), Inches(4.9)
    )
    image_panel.fill.solid()
    image_panel.fill.fore_color.rgb = COLORS["bg_soft"]
    image_panel.line.color.rgb = COLORS["primary"]
    image_panel.line.width = Pt(1)

    slide.shapes.add_picture(str(image_path), Inches(0.95), Inches(2.02), Inches(6.5), Inches(4.35))

    bullet_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(7.9), Inches(1.75), Inches(4.75), Inches(4.9)
    )
    bullet_panel.fill.solid()
    bullet_panel.fill.fore_color.rgb = COLORS["bg_soft"]
    bullet_panel.line.color.rgb = COLORS["accent"]
    bullet_panel.line.width = Pt(1)

    body_box = slide.shapes.add_textbox(Inches(8.2), Inches(2.05), Inches(4.2), Inches(4.35))
    body = body_box.text_frame
    body.clear()
    for idx, item in enumerate(bullets):
        p = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        p.font.name = "Noto Sans TC"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS["text"]

    return slide


def add_two_column_slide(prs, title, left_title, left_points, right_title, right_points):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    style_slide_background(slide, title)
    slide.shapes.title.text = title
    style_content_title(slide.shapes.title)

    left_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(1.75), Inches(5.95), Inches(4.9)
    )
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = COLORS["bg_soft"]
    left_card.line.color.rgb = COLORS["primary"]
    left_card.line.width = Pt(1)

    right_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.72), Inches(1.75), Inches(5.95), Inches(4.9)
    )
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = COLORS["bg_soft"]
    right_card.line.color.rgb = COLORS["accent"]
    right_card.line.width = Pt(1)

    left_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3))
    left_tf = left_box.text_frame
    left_tf.text = left_title
    left_tf.paragraphs[0].font.bold = True
    left_tf.paragraphs[0].font.name = "Noto Sans TC"
    left_tf.paragraphs[0].font.size = Pt(24)
    left_tf.paragraphs[0].font.color.rgb = COLORS["accent"]
    for point in left_points:
        p = left_tf.add_paragraph()
        p.text = f"• {point}"
        p.font.name = "Noto Sans TC"
        p.font.size = Pt(17)
        p.font.color.rgb = COLORS["text"]

    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.3))
    right_tf = right_box.text_frame
    right_tf.text = right_title
    right_tf.paragraphs[0].font.bold = True
    right_tf.paragraphs[0].font.name = "Noto Sans TC"
    right_tf.paragraphs[0].font.size = Pt(24)
    right_tf.paragraphs[0].font.color.rgb = COLORS["accent"]
    for point in right_points:
        p = right_tf.add_paragraph()
        p.text = f"• {point}"
        p.font.name = "Noto Sans TC"
        p.font.size = Pt(17)
        p.font.color.rgb = COLORS["text"]

    return slide


def collect_project_stats():
    root = Path("D:/AIClass")
    data_path = root / "js/data.js"
    notebook_path = root / "data/notebooklm-export.json"

    data_text = data_path.read_text(encoding="utf-8")
    node_ids = re.findall(r'id:\s*"([^"]+)"', data_text)
    years = [int(y) for y in re.findall(r"year:\s*(\d{4})", data_text)]
    categories = re.findall(r'category:\s*"([^"]+)"', data_text)
    youtube_ids = re.findall(r'youtubeId:\s*"([^"]*)"', data_text)

    node_count = len(node_ids)
    year_min = min(years) if years else 0
    year_max = max(years) if years else 0
    category_counts = Counter(categories)
    with_video = sum(1 for item in youtube_ids if item.strip())

    era_ranges = [
        ("1950s", 1950, 1959),
        ("1960-1989", 1960, 1989),
        ("1990-2009", 1990, 2009),
        ("2010-2019", 2010, 2019),
        ("2020-2026", 2020, 2026),
    ]
    era_counts = {}
    for label, start, end in era_ranges:
        era_counts[label] = sum(1 for year in years if start <= year <= end)

    notebook_nodes = []
    if notebook_path.exists():
        payload = json.loads(notebook_path.read_text(encoding="utf-8"))
        notebook_nodes = payload.get("nodes", [])

    notebook_ids = {item.get("nodeId") for item in notebook_nodes if item.get("nodeId")}
    notebook_matched = len(set(node_ids) & notebook_ids)

    source_files = [
        "index.html",
        "map.html",
        "topics.html",
        "about.html",
        "css/style.css",
        "css/map.css",
        "js/main.js",
        "js/map.js",
        "js/data.js",
    ]
    total_lines = 0
    for rel_path in source_files:
        file_path = root / rel_path
        if file_path.exists():
            total_lines += len(file_path.read_text(encoding="utf-8").splitlines())

    docs_for_counts = ["ReadMe.md", "index.html", "topics.html", "about.html", "task.md"]
    count_mentions = set()
    for rel_path in docs_for_counts:
        file_path = root / rel_path
        if not file_path.exists():
            continue
        text = file_path.read_text(encoding="utf-8")
        for match in re.findall(r"(\d+)\s*個?\s*節點", text):
            count_mentions.add(int(match))

    return {
        "node_count": node_count,
        "year_min": year_min,
        "year_max": year_max,
        "category_counts": category_counts,
        "with_video": with_video,
        "era_counts": era_counts,
        "notebook_count": len(notebook_nodes),
        "notebook_matched": notebook_matched,
        "total_lines": total_lines,
        "count_mentions": sorted(count_mentions),
    }


def main():
    prs = Presentation()
    stats = collect_project_stats()
    slides = []
    visuals = ensure_visual_assets()
    report_date = datetime.now().strftime("%Y/%m/%d")
    notebook_coverage = (stats["notebook_matched"] / stats["node_count"]) if stats["node_count"] else 0
    video_coverage = (stats["with_video"] / stats["node_count"]) if stats["node_count"] else 0
    category_text = "、".join(
        [
            f"應用突破 {stats['category_counts'].get('breakthrough', 0)}",
            f"模型架構 {stats['category_counts'].get('model', 0)}",
            f"理論基礎 {stats['category_counts'].get('theory', 0)}",
            f"訓練方法 {stats['category_counts'].get('training', 0)}",
            f"AI 寒冬 {stats['category_counts'].get('winter', 0)}",
        ]
    )
    era_text = "、".join(
        [
            f"1950s {stats['era_counts'].get('1950s', 0)}",
            f"1960-1989 {stats['era_counts'].get('1960-1989', 0)}",
            f"1990-2009 {stats['era_counts'].get('1990-2009', 0)}",
            f"2010-2019 {stats['era_counts'].get('2010-2019', 0)}",
            f"2020-2026 {stats['era_counts'].get('2020-2026', 0)}",
        ]
    )
    mention_text = " / ".join(str(x) for x in stats["count_mentions"]) if stats["count_mentions"] else "無"

    slides.append(add_title_slide(
        prs,
        "AI Explorer — 互動式 AI 地圖",
        f"期中成果報告\n班級：__________  姓名：__________  日期：{report_date}"
    ))

    slides.append(add_bullets_slide(prs, "目錄", [
        "1. 專案定位與解法",
        "2. 專案現況量化分析",
        "3. 核心功能與技術架構",
        "4. 風險盤點與期末規劃",
        "5. 結論與 Q&A",
    ]))

    slides.append(add_two_column_slide(
        prs,
        "專案定位與價值",
        "問題背景",
        ["AI 歷史與術語對新手門檻高", "單看文章難建立時間脈絡", "課堂展示缺乏可互動體驗"],
        "我們的解法",
        ["用地圖呈現 1950-2026 技術演進", "每個節點提供白話摘要 + 影片 + 小測驗", "整合 NotebookLM 來源，強化可查證性"],
    ))

    slides.append(add_bullets_slide(prs, "專案現況量化分析", [
        "網站頁面：4 頁（Home / Map / Topics / About）",
        f"核心程式碼：{stats['total_lines']} 行（HTML/CSS/JS）",
        f"資料節點：{stats['node_count']} 個，範圍 {stats['year_min']}–{stats['year_max']}",
        f"含影片節點：{stats['with_video']}/{stats['node_count']}（{video_coverage:.1%}）",
        f"NotebookLM 已同步：{stats['notebook_matched']}/{stats['node_count']}（{notebook_coverage:.1%}）",
    ]))

    slides.append(add_bullets_slide(prs, "資料層結構分析（js/data.js）", [
        f"類別分佈：{category_text}",
        f"年代分佈：{era_text}",
        "每個節點都含 id、年份、類別、敘述、誤解釐清、測驗",
        "連線欄位（connections）可視化技術影響路徑",
        "資料可被 map.js 與 topics 頁共用，維持一致資料源",
    ]))

    slides.append(add_bullets_slide(prs, "核心功能驗證（map.js）", [
        "年代篩選：全部 / 1950s / 1960-80s / 1990-2000s / 2010s / 2020s+",
        "互動操作：節點點擊、地圖拖曳、滑鼠滾輪縮放",
        "動態面板：摘要、影片、來源、常見誤解、小測驗即時回饋",
        "NotebookLM 覆蓋時優先顯示同步內容，否則回退預設資料",
        "主題索引頁與地圖共用資料，降低維護成本",
    ]))

    slides.append(add_image_bullets_slide(prs, "地圖截圖全覽", visuals["overview"], [
        "左側 65%：SVG 互動地圖",
        "右側 35%：節點詳情面板",
        "上方：年代篩選（1950s -> 2020s+）",
        "節點點擊後即時更新摘要、影片、來源、測驗",
    ]))

    slides.append(add_image_bullets_slide(prs, "節點 Demo 1 — 圖靈測試", visuals["turing"], [
        "年份：1950，類別：理論基礎",
        "目前唯一完成 NotebookLM 同步的節點",
        "本地 mp4 影片可直接在面板播放",
        "小測驗可即時顯示正確與錯誤回饋",
    ]))

    slides.append(add_image_bullets_slide(prs, "節點 Demo 2 — Transformer", visuals["transformer"], [
        "年份：2017，類別：模型架構",
        "定位：LLM 基礎架構，可串聯 BERT / GPT-3",
        "展示節點連線與脈絡化學習路徑",
        "可用年代篩選快速聚焦 2010s 深度學習階段",
    ]))

    slides.append(add_image_bullets_slide(prs, "節點 Demo 3 — ChatGPT", visuals["chatgpt"], [
        "年份：2022，類別：應用突破",
        "代表生成式 AI 大眾化轉折點",
        "可與 Multimodal / Agents 串成近代趨勢鏈",
        "目前來源覆蓋仍待補齊，適合作為期末強化重點",
    ]))

    slides.append(add_two_column_slide(
        prs,
        "技術架構與模組分工",
        "前端互動層",
        ["index/map/topics/about 四頁靜態網站", "map.js 負責地圖渲染、互動與面板更新", "CSS 設計系統 + RWD 斷點支援桌機與手機"],
        "資料與內容層",
        ["js/data.js 作為單一節點資料源", "notebooklm-export.json 作為可覆蓋內容來源", "scripts/ 內建 Python 腳本自動生成期中簡報"],
    ))

    slides.append(add_bullets_slide(prs, "期中完成度總結", [
        "功能完成：互動地圖、年代篩選、節點詳情、測驗回饋",
        "內容完成：33 節點資料結構與分類、連線、影片欄位",
        "頁面完成：首頁導覽、主題索引、關於頁與地圖核心頁",
        "整合完成：NotebookLM 資料載入與來源列表顯示機制",
        "簡報產出：以程式自動產生可重複更新的 PPT",
    ]))

    slides.append(add_bullets_slide(prs, "風險盤點與修正策略", [
        f"文件中的節點數字不一致：{mention_text}（需統一為 {stats['node_count']}）",
        "NotebookLM 同步覆蓋率偏低，內容品質尚未齊一",
        "2025-2026 節點帶有趨勢推估，需標示「預測性內容」",
        "task.md 內瀏覽器/RWD 驗收仍待完成，需補測試證據",
        "修正策略：先統一資料口徑，再做內容補齊與測試",
    ]))

    slides.append(add_bullets_slide(prs, "期末執行規劃（3 階段）", [
        "Phase 1（資料校正）：統一全站節點數與說明文案、補齊來源標註",
        "Phase 2（內容擴充）：將 NotebookLM 同步提升到 33/33",
        "Phase 3（體驗優化）：手機端互動與可讀性微調 + 使用者測試",
        "驗收指標：載入效能、節點正確率、測驗可用率、來源完整率",
        "交付成果：期末網站 Demo + 完整技術報告 + 最終簡報",
    ]))

    slides.append(add_bullets_slide(prs, "Q&A", [
        "結論：AI Explorer 已完成可互動的教學地圖雛形",
        "下一步聚焦內容完整度與資料一致性",
        "感謝聆聽，歡迎提問與建議",
    ]))

    for idx, slide in enumerate(prs.slides, start=1):
        add_page_number(slide, str(idx))

    output = Path("D:/AIClass/ppt/AI_Explorer_Midterm_Report.pptx")
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f"Generated: {output}")


if __name__ == "__main__":
    main()
